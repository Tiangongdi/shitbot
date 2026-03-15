#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT Generator Script
基于python-pptx库生成专业风格的PowerPoint演示文稿

使用方法:
    python create_ppt.py
    
或导入使用:
    from create_ppt import PPTGenerator
    generator = PPTGenerator()
    generator.add_title_slide("标题", "副标题")
    generator.save("output.pptx")
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class PPTGenerator:
    """PPT生成器类"""
    
    # 颜色方案
    COLORS = {
        'primary': RGBColor(31, 73, 125),      # 深蓝色
        'background': RGBColor(255, 255, 255), # 白色
        'text': RGBColor(50, 50, 50),          # 深灰色
        'subtitle': RGBColor(200, 200, 200),   # 浅灰色
        'accent': RGBColor(68, 114, 196),      # 强调蓝色
        'light_bg': RGBColor(240, 245, 250),   # 浅蓝背景
    }
    
    # 字体大小
    FONT_SIZES = {
        'title': 54,
        'page_title': 36,
        'content': 22,
        'two_column': 18,
        'subtitle': 28,
        'small': 16,
    }
    
    def __init__(self, title="演示文稿"):
        """初始化PPT生成器"""
        self.prs = Presentation()
        # 设置16:9宽屏比例
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.slide_count = 0
        
    def add_title_slide(self, title, subtitle="", date=None):
        """
        添加标题页（封面/结束页）
        
        参数:
            title: 主标题
            subtitle: 副标题（可选）
            date: 日期（可选，默认使用当前日期）
        """
        # 使用空白布局
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加深蓝色背景
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['primary']
        background.line.fill.background()
        
        # 添加主标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.FONT_SIZES['title'])
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(self.FONT_SIZES['subtitle'])
            p.font.color.rgb = self.COLORS['subtitle']
            p.alignment = PP_ALIGN.CENTER
        
        # 添加日期
        if date is None:
            date = datetime.now().strftime("%Y年%m月%d日")
        if date:
            date_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5),
                Inches(12.333), Inches(0.5)
            )
            date_frame = date_box.text_frame
            p = date_frame.paragraphs[0]
            p.text = date
            p.font.size = Pt(self.FONT_SIZES['small'])
            p.font.color.rgb = self.COLORS['subtitle']
            p.alignment = PP_ALIGN.CENTER
        
        self.slide_count += 1
        return slide
    
    def add_content_slide(self, title, content_list):
        """
        添加内容页
        
        参数:
            title: 页面标题
            content_list: 内容列表，每个元素是一个要点
                         支持二级列表，使用元组表示: ("主要点", ["子要点1", "子要点2"])
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题栏背景
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # 添加标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.FONT_SIZES['page_title'])
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 添加内容区域
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.6),
            Inches(11.933), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # 添加内容列表
        for i, item in enumerate(content_list):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # 处理二级列表
            if isinstance(item, tuple):
                main_point, sub_points = item
                p.text = f"• {main_point}"
                p.font.size = Pt(self.FONT_SIZES['content'])
                p.font.color.rgb = self.COLORS['text']
                p.space_after = Pt(6)
                
                for sub_item in sub_points:
                    p = content_frame.add_paragraph()
                    p.text = f"    ◦ {sub_item}"
                    p.font.size = Pt(self.FONT_SIZES['content'] - 2)
                    p.font.color.rgb = self.COLORS['text']
                    p.space_after = Pt(4)
            else:
                p.text = f"• {item}"
                p.font.size = Pt(self.FONT_SIZES['content'])
                p.font.color.rgb = self.COLORS['text']
                p.space_after = Pt(12)
        
        self.slide_count += 1
        return slide
    
    def add_two_column_slide(self, title, left_title, left_content, right_title, right_content):
        """
        添加双栏页
        
        参数:
            title: 页面标题
            left_title: 左栏标题
            left_content: 左栏内容列表
            right_title: 右栏标题
            right_content: 右栏内容列表
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题栏背景
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # 添加标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.FONT_SIZES['page_title'])
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 左栏标题
        left_title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(5.8), Inches(0.6)
        )
        left_title_frame = left_title_box.text_frame
        p = left_title_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(self.FONT_SIZES['subtitle'])
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # 左栏内容
        left_content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2),
            Inches(5.8), Inches(4.8)
        )
        left_content_frame = left_content_box.text_frame
        left_content_frame.word_wrap = True
        
        for i, item in enumerate(left_content):
            if i == 0:
                p = left_content_frame.paragraphs[0]
            else:
                p = left_content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(self.FONT_SIZES['two_column'])
            p.font.color.rgb = self.COLORS['text']
            p.space_after = Pt(8)
        
        # 右栏标题
        right_title_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.5),
            Inches(5.8), Inches(0.6)
        )
        right_title_frame = right_title_box.text_frame
        p = right_title_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(self.FONT_SIZES['subtitle'])
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # 右栏内容
        right_content_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(2.2),
            Inches(5.8), Inches(4.8)
        )
        right_content_frame = right_content_box.text_frame
        right_content_frame.word_wrap = True
        
        for i, item in enumerate(right_content):
            if i == 0:
                p = right_content_frame.paragraphs[0]
            else:
                p = right_content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(self.FONT_SIZES['two_column'])
            p.font.color.rgb = self.COLORS['text']
            p.space_after = Pt(8)
        
        # 添加分隔线
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.4), Inches(1.8),
            Inches(0.02), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.COLORS['subtitle']
        divider.line.fill.background()
        
        self.slide_count += 1
        return slide
    
    def add_section_slide(self, section_title, section_number=None):
        """
        添加章节分隔页
        
        参数:
            section_title: 章节标题
            section_number: 章节编号（可选）
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加浅蓝背景
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['light_bg']
        background.line.fill.background()
        
        # 添加章节编号
        if section_number:
            num_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5),
                Inches(12.333), Inches(1)
            )
            num_frame = num_box.text_frame
            p = num_frame.paragraphs[0]
            p.text = f"第{section_number}部分"
            p.font.size = Pt(self.FONT_SIZES['subtitle'])
            p.font.color.rgb = self.COLORS['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 添加章节标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(self.FONT_SIZES['title'] - 6)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        self.slide_count += 1
        return slide
    
    def add_image_slide(self, title, image_path, caption=""):
        """
        添加图片页
        
        参数:
            title: 页面标题
            image_path: 图片路径
            caption: 图片说明（可选）
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题栏背景
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # 添加标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.FONT_SIZES['page_title'])
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 添加图片
        if os.path.exists(image_path):
            # 计算图片位置和大小（居中显示）
            max_width = Inches(10)
            max_height = Inches(5)
            
            slide.shapes.add_picture(
                image_path,
                Inches(1.666), Inches(1.5),
                width=max_width, height=max_height
            )
        else:
            # 图片不存在时显示占位符
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.666), Inches(1.5),
                Inches(10), Inches(5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.COLORS['light_bg']
            
            text_box = slide.shapes.add_textbox(
                Inches(1.666), Inches(3.5),
                Inches(10), Inches(1)
            )
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"[图片未找到: {image_path}]"
            p.font.size = Pt(self.FONT_SIZES['content'])
            p.font.color.rgb = self.COLORS['text']
            p.alignment = PP_ALIGN.CENTER
        
        # 添加图片说明
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8),
                Inches(12.333), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            p = caption_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(self.FONT_SIZES['small'])
            p.font.color.rgb = self.COLORS['text']
            p.alignment = PP_ALIGN.CENTER
        
        self.slide_count += 1
        return slide
    
    def add_table_slide(self, title, headers, rows):
        """
        添加表格页
        
        参数:
            title: 页面标题
            headers: 表头列表
            rows: 数据行列表（每行是一个列表）
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题栏背景
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS['primary']
        title_bar.line.fill.background()
        
        # 添加标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.FONT_SIZES['page_title'])
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 添加表格
        num_rows = len(rows) + 1  # 包括表头
        num_cols = len(headers)
        
        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        ).table
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['primary']
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.size = Pt(self.FONT_SIZES['content'] - 2)
        
        # 设置数据行
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(self.FONT_SIZES['content'] - 2)
                paragraph.font.color.rgb = self.COLORS['text']
                
                # 交替行背景色
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.COLORS['light_bg']
        
        self.slide_count += 1
        return slide
    
    def save(self, output_path=None):
        """
        保存PPT文件
        
        参数:
            output_path: 输出路径（可选，默认保存到脚本所在目录）
        
        返回:
            保存的文件路径
        """
        if output_path is None:
            # 默认保存到脚本所在目录
            script_dir = os.path.dirname(os.path.abspath(__file__))
            filename = f"{self.title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path = os.path.join(script_dir, filename)
        
        # 确保目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        self.prs.save(output_path)
        print(f"PPT已保存到: {output_path}")
        print(f"共生成 {self.slide_count} 张幻灯片")
        return output_path


def create_demo_ppt():
    """创建演示PPT"""
    generator = PPTGenerator("演示文稿示例")
    
    # 封面
    generator.add_title_slide(
        "项目汇报演示",
        "ShitBot PPT Generator",
        "2026年3月14日"
    )
    
    # 目录页
    generator.add_content_slide(
        "目录",
        [
            "项目背景",
            "技术方案",
            "实施计划",
            "预期成果",
            "总结与展望"
        ]
    )
    
    # 内容页
    generator.add_content_slide(
        "项目背景",
        [
            "市场需求分析",
            "技术发展趋势",
            ("核心问题", [
                "用户痛点明确",
                "解决方案缺失",
                "市场机会巨大"
            ]),
            "项目定位与目标"
        ]
    )
    
    # 双栏页
    generator.add_two_column_slide(
        "技术方案对比",
        "方案A：传统架构",
        [
            "成熟稳定",
            "学习成本低",
            "社区支持完善",
            "扩展性一般"
        ],
        "方案B：微服务架构",
        [
            "高度可扩展",
            "技术栈现代",
            "运维复杂度高",
            "适合大规模应用"
        ]
    )
    
    # 表格页
    generator.add_table_slide(
        "实施计划",
        ["阶段", "时间", "主要任务", "负责人"],
        [
            ["第一阶段", "1-2月", "需求分析与设计", "张三"],
            ["第二阶段", "3-4月", "核心功能开发", "李四"],
            ["第三阶段", "5-6月", "测试与优化", "王五"],
            ["第四阶段", "7月", "部署上线", "赵六"]
        ]
    )
    
    # 章节分隔页
    generator.add_section_slide("总结与展望", "4")
    
    # 总结页
    generator.add_content_slide(
        "总结",
        [
            "项目目标明确，方案可行",
            "技术选型合理，风险可控",
            "团队配置完善，进度有保障",
            "预期成果显著，价值突出"
        ]
    )
    
    # 结束页
    generator.add_title_slide(
        "感谢聆听",
        "欢迎提问与交流"
    )
    
    return generator.save()


if __name__ == "__main__":
    print("=" * 50)
    print("PPT Generator - 演示文稿生成工具")
    print("=" * 50)
    print()
    
    # 创建演示PPT
    output_file = create_demo_ppt()
    
    print()
    print("=" * 50)
    print("生成完成！")
    print("=" * 50)
